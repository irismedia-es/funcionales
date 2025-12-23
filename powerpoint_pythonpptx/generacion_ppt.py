from pptx import Presentation
from pptx.util import Inches

def crear_presentacion():
    """
    Inicializa una nueva presentación de PowerPoint.
    
    Returns:
        Presentation: Objeto de presentación de pptx.
    """
    prs = Presentation()
    print("Nueva presentación creada.")
    return prs

def agregar_diapositiva_titulo(prs, titulo_texto, subtitulo_texto):
    """
    Agrega una diapositiva de título a la presentación.
    
    Args:
        prs (Presentation): La presentación actual.
        titulo_texto (str): El texto del título principal.
        subtitulo_texto (str): El texto del subtítulo.
        
    Returns:
        Slide: La diapositiva creada.
    """
    # Layout 0 suele ser "Title Slide"
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = titulo_texto
    subtitle.text = subtitulo_texto
    
    print(f"Diapositiva de título agregada: '{titulo_texto}'")
    return slide

def agregar_diapositiva_contenido(prs, titulo_texto, lista_puntos):
    """
    Agrega una diapositiva de contenido con viñetas.
    
    Args:
        prs (Presentation): La presentación actual.
        titulo_texto (str): Título de la diapositiva.
        lista_puntos (list): Lista de strings para los puntos del contenido.
        
    Returns:
        Slide: La diapositiva creada.
    """
    # Layout 1 suele ser "Title and Content"
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = titulo_texto
    
    # El placeholder del cuerpo suele ser el índice 1
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, punto in enumerate(lista_puntos):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = punto
        p.level = 0
        
    print(f"Diapositiva de contenido agregada: '{titulo_texto}'")
    return slide

def guardar_presentacion(prs, nombre_archivo):
    """
    Guarda la presentación en el disco.
    
    Args:
        prs (Presentation): Objeto de presentación.
        nombre_archivo (str): Ruta o nombre del archivo .pptx de salida.
    """
    prs.save(nombre_archivo)
    print(f"Presentación guardada exitosamente en: {nombre_archivo}")
