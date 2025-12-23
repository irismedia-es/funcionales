from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def crear_datos_grafico(categorias, series_dict):
    """
    Crea un objeto CategoryChartData a partir de listas de categorías y series.
    
    Args:
        categorias (list): Lista de nombres de categorías (eje X).
        series_dict (dict): Diccionario donde clave=nombre_serie y valor=lista_datos.
        
    Returns:
        CategoryChartData: Objeto de datos listo para el gráfico.
    """
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    
    for nombre_serie, valores in series_dict.items():
        chart_data.add_series(nombre_serie, valores)
        
    return chart_data

def agregar_grafico_barras(slide, chart_data, x, y, ancho, alto, titulo=None):
    """
    Agrega un gráfico de barras (columnas) a la diapositiva.
    
    Args:
        slide: Diapositiva destino.
        chart_data: Datos del gráfico.
        x, y, ancho, alto: Dimensiones y posición (usar Inches o Cm).
        titulo (str, opcional): Título del gráfico.
        
    Returns:
        GraphicFrame: El objeto gráfico insertado.
    """
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, ancho, alto, chart_data
    )
    chart = graphic_frame.chart
    
    if titulo:
        chart.has_title = True
        chart.chart_title.text_frame.text = titulo
        
    return graphic_frame

def agregar_grafico_torta(slide, chart_data, x, y, ancho, alto, titulo=None):
    """
    Agrega un gráfico de torta (Pie Chart) a la diapositiva.
    
    Args:
        slide: Diapositiva destino.
        chart_data: Datos del gráfico.
        x, y, ancho, alto: Dimensiones y posición.
        titulo (str, opcional): Título del gráfico.
        
    Returns:
        GraphicFrame: El objeto gráfico insertado.
    """
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, ancho, alto, chart_data
    )
    chart = graphic_frame.chart
    
    if titulo:
        chart.has_title = True
        chart.chart_title.text_frame.text = titulo
        
    # Mostrar etiquetas de datos y leyenda para tortas suele ser útil
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
        
    return graphic_frame
