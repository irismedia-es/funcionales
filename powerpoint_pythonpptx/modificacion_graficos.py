from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData

def buscar_grafico(slide):
    """
    Busca el primer gráfico presente en una diapositiva.
    
    Args:
        slide: Objeto diapositiva.
        
    Returns:
        Chart: El objeto chart encontrado o None.
    """
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return shape.chart
    return None

def actualizar_datos_grafico(chart, nuevas_categorias, nuevas_series_dict):
    """
    Reemplaza completamente los datos de un gráfico existente.
    
    Args:
        chart: Objeto Chart a modificar.
        nuevas_categorias (list): Lista de nuevas categorías.
        nuevas_series_dict (dict): Diccionario {nombre_serie: [valores]}.
    """
    try:
        from creacion_graficos import crear_datos_grafico
    except ImportError:
        # Fallback si se ejecuta desde otro contexto o si falla
        from .creacion_graficos import crear_datos_grafico
    
    chart_data = crear_datos_grafico(nuevas_categorias, nuevas_series_dict)
    chart.replace_data(chart_data)
    print("Datos del gráfico actualizados.")

def modificar_titulo_grafico(chart, nuevo_titulo):
    """
    Cambia el título de un gráfico.
    """
    if not chart.has_title:
        chart.has_title = True
        
    chart.chart_title.text_frame.text = nuevo_titulo
