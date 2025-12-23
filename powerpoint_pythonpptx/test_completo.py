import os
from pptx.util import Inches

# Importar nuestros módulos
import generacion_ppt
import lectura_ppt
import analisis_diapositiva
import reemplazo_elementos
import creacion_graficos
import modificacion_graficos

def main():
    print("=== INICIO DE PRUEBA ===")
    
    # 1. Crear Presentación
    print("\n1. Generando PPT Base...")
    prs = generacion_ppt.crear_presentacion()
    generacion_ppt.agregar_diapositiva_titulo(prs, "Prueba Automatizada", "Generado con Python")
    slide_contenido = generacion_ppt.agregar_diapositiva_contenido(prs, "Agenda", ["Punto A", "Punto B", "Punto C"])
    
    # 2. Agregar Gráfico
    print("\n2. Agregando Gráfico...")
    slide_grafico = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide_grafico.shapes.title.text = "Ventas Mensuales"
    
    datos_cats = ['Ene', 'Feb', 'Mar']
    datos_series = {'Ventas': [10, 20, 30]}
    chart_data = creacion_graficos.crear_datos_grafico(datos_cats, datos_series)
    
    creacion_graficos.agregar_grafico_barras(slide_grafico, chart_data, Inches(1), Inches(2), Inches(6), Inches(4))
    
    nombre_temp = "temp_prueba.pptx"
    generacion_ppt.guardar_presentacion(prs, nombre_temp)
    
    # 3. Leer y Modificar
    print("\n3. Leyendo y Modificando...")
    prs_leida = lectura_ppt.cargar_presentacion(nombre_temp)
    
    # Analizar slide 1 (Agenda)
    analisis_diapositiva.imprimir_estructura_diapositiva(prs_leida.slides[1])
    
    # Reemplazar texto en slide 1
    slide_agenda = prs_leida.slides[1]
    reemplazo_elementos.reemplazar_texto_en_diapositiva(slide_agenda, "Punto A", "Punto Alpha Modificado")
    
    # Modificar gráfico en slide 2
    slide_graf_leida = prs_leida.slides[2]
    chart = modificacion_graficos.buscar_grafico(slide_graf_leida)
    if chart:
        print("Gráfico encontrado, actualizando datos...")
        nuevos_datos = {'Ventas': [50, 60, 100]} # Cambian los valores
        modificacion_graficos.actualizar_datos_grafico(chart, datos_cats, nuevos_datos)
        modificacion_graficos.modificar_titulo_grafico(chart, "Ventas Actualizadas (Q1)")
    
    # 4. Guardar Final
    nombre_final = "resultado_final.pptx"
    generacion_ppt.guardar_presentacion(prs_leida, nombre_final)
    
    print(f"\n=== PRUEBA COMPLETADA ===")
    print(f"Archivos generados: {nombre_temp}, {nombre_final}")

if __name__ == "__main__":
    main()
